"""PPTX parser — slide-aware section splitting.

Uses python-pptx for slide extraction.
Extracts title, bullet text, speaker notes, and shape text.
Merges 1-3 slides per section.
"""

from __future__ import annotations

from pathlib import Path

from rag.source_section import SourceSection


def parse(path: str | Path, base_dir: str | Path = ".") -> list[SourceSection]:
    """Parse a PPTX file into slide-aware SourceSections."""
    path = Path(path)
    base_dir = Path(base_dir)

    try:
        relative = path.relative_to(base_dir)
    except ValueError:
        relative = path

    source = str(relative).replace("\\", "/")
    doc_title = path.stem

    try:
        from pptx import Presentation
    except ImportError:
        return []

    try:
        prs = Presentation(str(path))
    except Exception:
        return []

    slides: list[dict] = []
    for i, slide in enumerate(prs.slides):
        slide_data = _extract_slide(slide, i + 1)
        slides.append(slide_data)

    if not slides:
        return []

    # Group slides into sections (1-3 slides per section)
    return _group_slides(slides, source, doc_title)


def _extract_slide(slide, slide_num: int) -> dict:
    """Extract text content from a single slide."""
    title = ""
    body_parts = []
    notes = ""
    image_count = 0

    for shape in slide.shapes:
        # Count embedded pictures (P5G.0: image-only slides are reported,
        # not silently dropped).
        try:
            if shape.shape_type is not None and str(shape.shape_type) == "PICTURE (13)":
                image_count += 1
            elif getattr(shape, "image", None) is not None:
                image_count += 1
        except Exception:
            pass

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # First text with title placeholder or large font is the title
                if not title and (
                    shape.shape_type is not None and
                    hasattr(shape, "placeholder_format") and
                    shape.placeholder_format is not None and
                    shape.placeholder_format.idx == 0
                ):
                    title = text
                else:
                    body_parts.append(text)

        # Speaker notes
        if shape.has_text_frame and hasattr(slide, "notes_slide"):
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

    # Also try to get notes from the slide's notes_slide directly
    if not notes:
        try:
            ns = slide.notes_slide
            if ns and ns.notes_text_frame:
                notes = ns.notes_text_frame.text.strip()
        except Exception:
            pass

    body = "\n".join(body_parts)
    if notes:
        body += f"\n\n[Notes: {notes}]"

    return {
        "slide_num": slide_num,
        "title": title,
        "text": body,
        "image_count": image_count,
        "has_text": bool(title or body_parts or notes),
    }


def _group_slides(
    slides: list[dict], source: str, doc_title: str, max_slides: int = 3
) -> list[SourceSection]:
    """Group slides into sections for better chunk context."""
    sections: list[SourceSection] = []
    buffer: list[dict] = []

    for slide in slides:
        buffer.append(slide)
        if len(buffer) >= max_slides:
            sections.append(_make_section(buffer, source, doc_title))
            buffer = []

    if buffer:
        sections.append(_make_section(buffer, source, doc_title))

    return sections


def _make_section(slides: list[dict], source: str, doc_title: str) -> SourceSection:
    """Create a SourceSection from a group of slides."""
    first = slides[0]
    last = slides[-1]

    parts = []
    for s in slides:
        if s["title"]:
            parts.append(f"## {s['title']}")
        if s["text"]:
            parts.append(s["text"])

    heading_path = [first["title"]] if first["title"] else []

    return SourceSection(
        source=source,
        doc_title=doc_title,
        unit_type="slide",
        unit_range=f"slide {first['slide_num']}-{last['slide_num']}",
        heading_path=heading_path,
        text="\n\n".join(parts),
        metadata={
            "file_type": "pptx",
            "slide_start": first["slide_num"],
            "slide_end": last["slide_num"],
            "image_count": sum(s["image_count"] for s in slides),
            "empty_slides": sum(0 if s["has_text"] else 1 for s in slides),
        },
    )
