"""Tests for rag/parsers/report.py — extraction reports (P5G.0)."""

from pathlib import Path

import pytest

from rag.parsers import parse_document
from rag.parsers.report import build_report, error_report, unit_kind_for
from rag.source_section import SourceSection


def _section(text: str, *, file_type: str = "pdf", **meta) -> SourceSection:
    metadata = {"file_type": file_type}
    metadata.update(meta)
    return SourceSection(
        source="test.pdf",
        doc_title="test",
        unit_type="page",
        unit_range="p1",
        heading_path=[],
        text=text,
        metadata=metadata,
    )


class TestBuildReport:
    def test_complete_when_all_units_have_text(self):
        report = build_report("pdf", "rag.parsers.pdf_parser", [
            _section("page one"),
            _section("page two"),
        ])
        assert report.status == "complete"
        assert report.total_units == 2
        assert report.extracted_units == 2
        assert report.empty_units == 0
        assert report.warnings == []

    def test_partial_when_some_pages_empty(self):
        report = build_report("pdf", "rag.parsers.pdf_parser", [
            _section("page one"),
            _section("", needs_ocr=True),
            _section("page three"),
        ])
        assert report.status == "partial"
        assert report.total_units == 3
        assert report.extracted_units == 2
        assert report.empty_units == 1
        assert "scanned_or_empty_pages" in report.warnings

    def test_empty_when_no_text_at_all(self):
        report = build_report("pptx", "rag.parsers.pptx_parser", [
            _section("", file_type="pptx"),
            _section("", file_type="pptx"),
        ])
        assert report.status == "empty"
        assert "slides_without_text" in report.warnings
        assert "no_searchable_text" in report.warnings

    def test_error_report(self):
        report = error_report("pdf", "rag.parsers.pdf_parser")
        assert report.status == "error"
        assert report.warnings == ["parser_error"]

    def test_image_only_document_warns(self):
        report = build_report("docx", "rag.parsers.docx_parser", [], image_count=5)
        assert report.status == "empty"
        assert "images_not_recognized" in report.warnings

    def test_unit_kind_per_format(self):
        assert unit_kind_for("pdf") == "page"
        assert unit_kind_for("pptx") == "slide"
        assert unit_kind_for("docx") == "heading section"

    def test_to_dict_shape(self):
        data = build_report("pdf", "rag.parsers.pdf_parser", [_section("x")]).to_dict()
        assert data["status"] == "complete"
        assert data["total_units"] == 1
        assert data["parser"] == "rag.parsers.pdf_parser"
        assert data["file_type"] == "pdf"


class TestParseDocumentReport:
    def test_unsupported_extension_returns_error_report(self, tmp_path):
        f = tmp_path / "notes.txt"
        f.write_text("not a typed format here", encoding="utf-8")
        # .txt routes to markdown_parser, so use an unknown extension instead
        f = tmp_path / "notes.xyz"
        f.write_text("x", encoding="utf-8")
        sections, report = parse_document(f, tmp_path)
        assert sections == []
        assert report.status == "error"

    def test_markdown_returns_complete_report(self, tmp_path):
        f = tmp_path / "note.md"
        f.write_text("# Title\n\nBody text here.\n", encoding="utf-8")
        sections, report = parse_document(f, tmp_path)
        assert sections
        assert report.status == "complete"
        assert report.total_units == len(sections)


def _write_two_page_pdf(path: Path) -> None:
    """Build a tiny 2-page text PDF with pypdf (no external tools)."""
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    font = DictionaryObject({NameObject("/F1"): DictionaryObject({
        NameObject("/Type"): NameObject("/Font"),
        NameObject("/Subtype"): NameObject("/Type1"),
        NameObject("/BaseFont"): NameObject("/Helvetica"),
    })})
    for label in (b"Hello Bobodan page one", b"Second page content"):
        stream = DecodedStreamObject()
        stream.set_data(b"BT /F1 12 Tf 72 720 Td (" + label + b") Tj ET")
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Contents")] = writer._add_object(stream)
        page[NameObject("/Resources")] = writer._add_object(
            DictionaryObject({NameObject("/Font"): font})
        )
        writer.add_page(page)
    with open(path, "wb") as handle:
        writer.write(handle)


class TestPdfParserPypdfOnly:
    def test_parses_text_pages(self, tmp_path):
        pdf = tmp_path / "doc.pdf"
        _write_two_page_pdf(pdf)
        sections, report = parse_document(pdf, tmp_path)
        assert sections, "pypdf should extract text from both pages"
        assert report.status == "complete"
        assert report.total_units == 2
        assert report.parser == "rag.parsers.pdf_parser"
        assert report.file_type == "pdf"

    def test_parser_module_has_no_fitz_reference(self):
        import rag.parsers.pdf_parser as pdf_parser
        source = Path(pdf_parser.__file__).read_text(encoding="utf-8")
        assert "fitz" not in source
        assert "PyMuPDF" not in source

    def test_non_pdf_content_returns_error_report(self, tmp_path):
        # Missing the %PDF- magic header → parser treats it as a parse failure.
        pdf = tmp_path / "broken.pdf"
        pdf.write_bytes(b"this is not a pdf at all")
        sections, report = parse_document(pdf, tmp_path)
        assert sections == []
        assert report.status == "error"
        assert report.warnings == ["parser_error"]

    def test_image_only_pdf_is_empty_not_error(self, tmp_path):
        # A structurally valid PDF with no text layer is "empty" (scanned),
        # not a parser failure — the Library must show it as not retrievable.
        from pypdf import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        pdf = tmp_path / "scanned.pdf"
        with open(pdf, "wb") as handle:
            writer.write(handle)
        sections, report = parse_document(pdf, tmp_path)
        assert len(sections) == 1
        assert sections[0].metadata.get("needs_ocr") is True
        assert report.status == "empty"
        assert "scanned_or_empty_pages" in report.warnings


# ── DOCX / PPTX image statistics (P5G.0) ───────────────────────────────

_1PX_PNG = (
    b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _png_bytes() -> bytes:
    import base64
    return base64.b64decode(_1PX_PNG)


class TestDocxImageStats:
    def test_docx_with_inline_image_reports_image_count(self, tmp_path):
        from docx import Document

        img = tmp_path / "pixel.png"
        img.write_bytes(_png_bytes())
        doc = Document()
        doc.add_heading("Title", level=1)
        doc.add_paragraph("Body text with an image below.")
        doc.add_picture(str(img))
        doc.save(str(tmp_path / "with-img.docx"))

        sections, report = parse_document(tmp_path / "with-img.docx", tmp_path)
        assert report.status == "complete"
        assert report.image_count == 1
        assert sum(s.metadata.get("image_count") or 0 for s in sections) == 1

    def test_image_only_docx_paragraph_still_counts(self, tmp_path):
        from docx import Document

        img = tmp_path / "pixel.png"
        img.write_bytes(_png_bytes())
        doc = Document()
        doc.add_heading("Figures", level=1)
        doc.add_picture(str(img))
        doc.add_picture(str(img))
        doc.save(str(tmp_path / "figures.docx"))

        sections, report = parse_document(tmp_path / "figures.docx", tmp_path)
        assert report.image_count == 2
        assert "images_not_recognized" not in report.warnings  # has some text? no —
        # image-only section text is empty, so expect no_searchable_text warning
        assert report.status in {"empty", "partial"}


class TestPptxImageStats:
    def test_pptx_picture_slide_reports_image_and_warning(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches

        img = tmp_path / "pixel.png"
        img.write_bytes(_png_bytes())
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        slide1.shapes.title.text = "Slide One"
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        slide2.shapes.add_picture(str(img), Inches(1), Inches(1))
        prs.save(str(tmp_path / "mixed.pptx"))

        sections, report = parse_document(tmp_path / "mixed.pptx", tmp_path)
        assert report.image_count == 1
        assert "slides_without_text" in report.warnings
